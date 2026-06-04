#!/usr/bin/env python3
"""
wiki-upload-extract.py — Extract text from uploaded documents.

Supports:
  - .pages (Apple Pages) via IWA/snappy protobuf parsing
  - .pdf  via MinerU-first deep extraction, with explicit degraded fallback
  - .docx/.pptx via document-explorer dependencies when available
  - .html/.htm via HTMLParser
  - .md/.txt pass-through

Writes derived text artifacts to <upload_dir>/_extracted/.
Emits explicit extract_failed records when extraction fails.

Usage:
  python3 wiki-upload-extract.py --source <file_or_dir> [--batch <prefix>] [--json]
"""

import argparse
import json
import os
import re
import subprocess
import sys
import zipfile
from pathlib import Path


VAULT_ROOT = Path(os.environ.get("OBSIDIAN_VAULT_PATH", str(Path.home() / "Knowledge"))).expanduser()
MINERU_EXTRACTOR = os.environ.get("SOLAR_HARNESS_MINERU_EXTRACTOR", "")
REQUIRE_MINERU = os.environ.get("SOLAR_HARNESS_REQUIRE_MINERU", "0") == "1"
MAX_MINERU_ARTIFACT_CHARS = int(os.environ.get("SOLAR_HARNESS_MINERU_ARTIFACT_CHARS", "240000"))


def _base_result() -> dict:
    return {
        "status": "pending",
        "text": "",
        "error": None,
        "method": None,
        "quality": "unknown",
        "provenance": {},
    }


def _read_limited(path: Path, max_chars: int) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")[:max_chars]
    except OSError:
        return ""


def _mineru_command(pdf_path: str) -> list[str]:
    if MINERU_EXTRACTOR:
        return [MINERU_EXTRACTOR, pdf_path, str(VAULT_ROOT)]
    return [
        sys.executable,
        str(Path(__file__).resolve().with_name("mineru_extract.py")),
        "extract",
        pdf_path,
        "--vault",
        str(VAULT_ROOT),
        "--json",
    ]


def _collect_mineru_markdown(payload: dict) -> str:
    ref_dir = Path(str(payload.get("ref_dir", ""))).expanduser()
    pages = payload.get("generated_pages") or []
    candidates: list[Path] = []
    if ref_dir:
        candidates.append(ref_dir / "index.md")
    for item in pages:
        p = Path(str(item)).expanduser()
        if not p.is_absolute() and ref_dir:
            p = ref_dir / p
        candidates.append(p)

    seen: set[str] = set()
    parts: list[str] = []
    remaining = MAX_MINERU_ARTIFACT_CHARS
    for path in candidates:
        key = str(path)
        if key in seen or remaining <= 0:
            continue
        seen.add(key)
        text = _read_limited(path, remaining)
        if not text.strip():
            continue
        parts.append(f"\n\n<!-- mineru_source: {path} -->\n\n{text}")
        remaining -= len(text)
    return "\n".join(parts).strip()

# ---------------------------------------------------------------------------
# Pages extraction (IWA + snappy)
# ---------------------------------------------------------------------------

def extract_pages_text(pages_path: str) -> dict:
    """Extract text from an Apple .pages file by parsing IWA archives."""
    result = _base_result()

    if not os.path.isfile(pages_path):
        result["status"] = "extract_failed"
        result["error"] = "file_not_found"
        return result

    # Method 1: IWA + snappy decompression (works for modern .pages zip format)
    try:
        import snappy
        text_parts = []
        with zipfile.ZipFile(pages_path) as zf:
            iwa_files = [n for n in zf.namelist() if n.endswith('.iwa')]
            for iwa in iwa_files:
                try:
                    with zf.open(iwa) as f:
                        data = f.read()
                    if len(data) < 5:
                        continue
                    decompressed = snappy.decompress(data[4:])
                    raw = decompressed.decode('utf-8', errors='replace')
                    # Chinese + mixed text
                    chinese = re.findall(
                        r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+'
                        r'(?:[a-zA-Z0-9\s\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]*)*',
                        raw,
                    )
                    # Long ASCII runs
                    ascii_runs = re.findall(r'[a-zA-Z][a-zA-Z0-9\s,.\-:;!?()]{15,}', raw)
                    text_parts.extend(chinese)
                    text_parts.extend(ascii_runs)
                except Exception:
                    pass

        if text_parts:
            result["text"] = '\n'.join(text_parts)
            result["status"] = "success"
            result["method"] = "iwa_snappy"
            result["quality"] = "structured_text"
            return result
    except ImportError:
        pass
    except zipfile.BadZipFile:
        pass

    # Method 2: qlmanage HTML preview fallback (macOS only)
    try:
        import subprocess
        import tempfile
        with tempfile.TemporaryDirectory() as tmpdir:
            proc = subprocess.run(
                ['qlmanage', '-p', '-o', tmpdir, pages_path],
                capture_output=True, text=True, timeout=30,
            )
            # Find Preview.html in the qlpreview directory
            for d in os.listdir(tmpdir):
                ql_dir = os.path.join(tmpdir, d)
                if os.path.isdir(ql_dir):
                    preview = os.path.join(ql_dir, 'Preview.html')
                    if os.path.isfile(preview):
                        with open(preview) as f:
                            html = f.read()
                        from html.parser import HTMLParser
                        class _TE(HTMLParser):
                            def __init__(self):
                                super().__init__()
                                self.parts = []
                                self._skip = False
                            def handle_starttag(self, tag, attrs):
                                if tag in ('script', 'style'):
                                    self._skip = True
                            def handle_endtag(self, tag):
                                if tag in ('script', 'style'):
                                    self._skip = False
                            def handle_data(self, data):
                                if not self._skip:
                                    t = data.strip()
                                    if t:
                                        self.parts.append(t)
                        te = _TE()
                        te.feed(html)
                        if te.parts:
                            result["text"] = '\n'.join(te.parts)
                            result["status"] = "success"
                            result["method"] = "qlmanage_html"
                            result["quality"] = "fallback_preview"
                            return result
    except Exception:
        pass

    result["status"] = "extract_failed"
    result["error"] = "all_methods_exhausted"
    return result


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------

def extract_pdf_text(pdf_path: str) -> dict:
    """Extract PDF through MinerU first.

    PyMuPDF is retained only as an explicit degraded fallback so upload/backfill
    paths cannot silently create low-quality paper notes from raw PDF text.
    """
    result = _base_result()
    mineru_error = ""
    try:
        proc = subprocess.run(
            _mineru_command(pdf_path),
            capture_output=True,
            text=True,
            timeout=int(os.environ.get("SOLAR_HARNESS_MINERU_TIMEOUT", "900")),
        )
        if proc.returncode == 0:
            payload = json.loads(proc.stdout or "{}")
            text = _collect_mineru_markdown(payload)
            if payload.get("ok") and text.strip():
                result["text"] = text
                result["status"] = "success"
                result["method"] = f"mineru:{payload.get('method') or 'unknown'}"
                result["quality"] = "mineru_deep_extraction"
                result["provenance"] = {
                    "mineru_ref_dir": payload.get("ref_dir", ""),
                    "mineru_generated_pages": payload.get("generated_pages", []),
                    "mineru_extraction_status": payload.get("extraction_status", ""),
                }
                return result
            mineru_error = payload.get("error") or payload.get("extraction_status") or "mineru_empty_artifact"
        else:
            mineru_error = (proc.stderr or proc.stdout or f"mineru_exit_{proc.returncode}").strip()[:1000]
    except Exception as e:
        mineru_error = f"{type(e).__name__}: {e}"

    if REQUIRE_MINERU:
        result["status"] = "extract_failed"
        result["error"] = f"mineru_required:{mineru_error or 'unavailable'}"
        result["method"] = "mineru_required"
        result["quality"] = "blocked"
        return result

    try:
        import fitz
        doc = fitz.open(pdf_path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        text = '\n'.join(pages)
        if text.strip():
            result["text"] = text
            result["status"] = "success"
            result["method"] = "mineru_unavailable:pymupdf_fallback"
            result["quality"] = "degraded_fallback_requires_reingest"
            result["provenance"] = {"mineru_error": mineru_error}
        else:
            result["status"] = "extract_failed"
            result["error"] = f"mineru_failed_then_pymupdf_empty:{mineru_error}"
    except ImportError:
        result["status"] = "extract_failed"
        result["error"] = f"mineru_failed_and_pymupdf_not_installed:{mineru_error}"
    except Exception as e:
        result["status"] = "extract_failed"
        result["error"] = f"mineru_failed_then_pymupdf_error:{mineru_error}; {e}"
    return result


def extract_docx_text(docx_path: str) -> dict:
    result = _base_result()
    try:
        import docx  # type: ignore
        document = docx.Document(docx_path)
        parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        if text.strip():
            result["text"] = text
            result["status"] = "success"
            result["method"] = "document_explorer:python_docx"
            result["quality"] = "structured_text"
        else:
            result["status"] = "extract_failed"
            result["error"] = "empty_text"
            result["quality"] = "blocked"
    except ImportError:
        result["status"] = "extract_failed"
        result["error"] = "document_explorer_dependency_missing:python-docx"
        result["quality"] = "blocked"
    except Exception as e:
        result["status"] = "extract_failed"
        result["error"] = str(e)
        result["quality"] = "blocked"
    return result


def extract_pptx_text(pptx_path: str) -> dict:
    result = _base_result()
    try:
        from pptx import Presentation  # type: ignore
        presentation = Presentation(pptx_path)
        parts: list[str] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_parts.append(text.strip())
            if slide_parts:
                parts.append(f"## Slide {idx}\n" + "\n".join(slide_parts))
        text = "\n\n".join(parts)
        if text.strip():
            result["text"] = text
            result["status"] = "success"
            result["method"] = "document_explorer:python_pptx"
            result["quality"] = "structured_text"
        else:
            result["status"] = "extract_failed"
            result["error"] = "empty_text"
            result["quality"] = "blocked"
    except ImportError:
        result["status"] = "extract_failed"
        result["error"] = "document_explorer_dependency_missing:python-pptx"
        result["quality"] = "blocked"
    except Exception as e:
        result["status"] = "extract_failed"
        result["error"] = str(e)
        result["quality"] = "blocked"
    return result


# ---------------------------------------------------------------------------
# HTML extraction
# ---------------------------------------------------------------------------

def extract_html_text(html_path: str) -> dict:
    from html.parser import HTMLParser
    result = _base_result()
    try:
        with open(html_path, encoding='utf-8', errors='replace') as f:
            html = f.read()
        class _TE(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ('script', 'style'):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ('script', 'style'):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    t = data.strip()
                    if t:
                        self.parts.append(t)
        te = _TE()
        te.feed(html)
        text = '\n'.join(te.parts)
        if text.strip():
            result["text"] = text
            result["status"] = "success"
            result["method"] = "html_parser"
            result["quality"] = "structured_text"
        else:
            result["status"] = "extract_failed"
            result["error"] = "empty_text"
    except Exception as e:
        result["status"] = "extract_failed"
        result["error"] = str(e)
    return result


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

EXTRACTORS = {
    '.pages': extract_pages_text,
    '.pdf': extract_pdf_text,
    '.docx': extract_docx_text,
    '.pptx': extract_pptx_text,
    '.html': extract_html_text,
    '.htm': extract_html_text,
}


def extract_file(source_path: str, output_dir: str = None) -> dict:
    """Extract text from a single file. Returns result dict with artifact path."""
    _, ext = os.path.splitext(source_path)
    ext = ext.lower()

    # Pass-through for already-text formats
    if ext in ('.md', '.txt', '.markdown'):
        with open(source_path, encoding='utf-8', errors='replace') as f:
            text = f.read()
        return {
            "source": os.path.basename(source_path),
            "status": "success",
            "method": "passthrough",
            "quality": "source_text",
            "error": None,
            "text_length": len(text),
            "text_preview": text[:500],
            "artifact": source_path,
        }

    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return {
            "source": os.path.basename(source_path),
            "status": "extract_failed",
            "method": None,
            "quality": "blocked",
            "error": f"unsupported_extension:{ext}",
            "text_length": 0,
            "text_preview": None,
            "artifact": None,
        }

    res = extractor(source_path)
    res["source"] = os.path.basename(source_path)

    # Write derived artifact
    if output_dir is None:
        output_dir = os.path.join(os.path.dirname(source_path), '_extracted')
    os.makedirs(output_dir, exist_ok=True)

    basename = os.path.basename(source_path)
    name_without_ext = os.path.splitext(basename)[0]
    artifact_path = os.path.join(output_dir, f"{name_without_ext}.extracted.txt")

    with open(artifact_path, 'w', encoding='utf-8') as f:
        if res["status"] == "success":
            f.write(res["text"])
        else:
            f.write(f"[extract_failed: {res['error']}]")

    res["artifact"] = artifact_path
    res["text_length"] = len(res.get("text", ""))
    res["text_preview"] = res.get("text", "")[:500]
    # Don't keep full text in the summary
    res.pop("text", None)

    return res


def extract_batch(upload_dir: str, batch_prefix: str = None) -> dict:
    """Extract all files matching batch_prefix in upload_dir."""
    files = sorted(os.listdir(upload_dir))
    if batch_prefix:
        files = [f for f in files if f.startswith(batch_prefix) and not f.startswith('.')]

    results = []
    for f in files:
        path = os.path.join(upload_dir, f)
        if not os.path.isfile(path):
            continue
        res = extract_file(path)
        results.append(res)

    summary = {
        "batch": batch_prefix or "all",
        "upload_dir": upload_dir,
        "total_files": len(results),
        "extracted": sum(1 for r in results if r["status"] == "success"),
        "failed": sum(1 for r in results if r["status"] == "extract_failed"),
        "files": results,
    }
    return summary


def main():
    parser = argparse.ArgumentParser(description='Extract text from uploaded documents')
    parser.add_argument('--source', help='Single file or directory to extract')
    parser.add_argument('--batch', help='Batch prefix filter (e.g., 20260508T122047Z)')
    parser.add_argument('--json', action='store_true', help='Output JSON')
    parser.add_argument('--output-dir', help='Override output directory for derived artifacts')
    args = parser.parse_args()

    if not args.source:
        print("Error: --source required", file=sys.stderr)
        sys.exit(1)

    if os.path.isdir(args.source):
        summary = extract_batch(args.source, args.batch)
        if args.json:
            print(json.dumps(summary, indent=2, ensure_ascii=False))
        else:
            print(f"Batch: {summary['batch']}")
            print(f"Total: {summary['total_files']}, Extracted: {summary['extracted']}, Failed: {summary['failed']}")
            for f in summary['files']:
                icon = '✓' if f['status'] == 'success' else '✗'
                print(f"  {icon} {f['source']}: {f['status']} ({f.get('text_length', 0)} chars)")
    else:
        res = extract_file(args.source, args.output_dir)
        if args.json:
            print(json.dumps(res, indent=2, ensure_ascii=False))
        else:
            print(f"{res['source']}: {res['status']} ({res.get('text_length', 0)} chars)")


if __name__ == '__main__':
    main()
